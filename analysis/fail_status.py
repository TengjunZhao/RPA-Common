import pandas as pd
import numpy as np
from sqlalchemy import create_engine, text
from datetime import datetime, timedelta
import calendar
from sqlalchemy.exc import SQLAlchemyError
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from decimal import Decimal
from pptx.util import Pt



def get_first_and_last_day_of_month(month_str):
    """
    Get the first and last day of the month given a string in YYYYMM format.
    :param month_str: str, the month in YYYYMM format.
    :return: tuple, (first_day, last_day) in YYYY-MM-DD format.
    """
    year = int(month_str[:4])
    month = int(month_str[4:])

    # First day of the month
    first_day = f"{year}{month:02d}01"

    # Last day of the month
    last_day = f"{year}{month:02d}{calendar.monthrange(year, month)[1]}"

    return first_day, last_day

def fetch_data_from_db(connection, query):
    """Fetch data from the database using the provided query."""
    return pd.read_sql(query, connection)

def process_data_from_db(db_config, month_str):
    """Process data from the database and compute the required metrics."""
    first_day, last_day = get_first_and_last_day_of_month(month_str)
    try:
        # Establish SQLAlchemy connection
        engine = create_engine(
            f"mysql+pymysql://{db_config['user']}:"
            f"{db_config['password']}@{db_config['host']}:"
            f"{db_config['port']}/{db_config['database']}"
        )

        # Query to fetch required data
        query = f"""
        SELECT 
            device,
            oper_old, 
            SUM(in_qty) AS m_In, 
            SUM(out_qty) AS m_Out 
        FROM 
            cmsalpha.db_yielddetail 
        WHERE 
            workdt BETWEEN '{first_day}' AND '{last_day}'
        GROUP BY 
            device, oper_old
        """
        # Fetch data
        data = fetch_data_from_db(engine, query)

        # Pivot table for calculations
        pivot_df = data.pivot_table(index='device', columns='oper_old', values=['m_In', 'm_Out'], aggfunc='sum')
        pivot_df = pivot_df.fillna(0)
        valid_operators = ['5700', '5710', '5780']
        for operator in valid_operators:
            if operator not in pivot_df.columns.get_level_values(1):  # Check if the operator exists in columns
                pivot_df[('m_In', operator)] = np.nan  # Add missing operator column with NaN values
                pivot_df[('m_Out', operator)] = np.nan  # Add missing operator column with NaN values
                # Calculate Yield
        yield_df = pivot_df['m_Out'] / pivot_df['m_In']
        # 如果 yield_df 的列索引是多级索引，需要通过元组形式访问列
        # Calculate Yield_5700_5710_5780, ignoring NaN and 0 values
        def calculate_yield(row):
            # Filter out NaN and 0 values from the row
            valid_values = row[row != 0].dropna()
            if len(valid_values) > 0:
                return valid_values.prod()  # Multiply all valid values together
            else:
                return np.nan  # If no valid values, return NaN

        # Apply the function to calculate the Yield for each row
        yield_df['Yield_5700_5710_5780'] = yield_df[valid_operators].apply(calculate_yield, axis=1)
        # Calculate input quantities
        input_df = pivot_df['m_In'][['5700', '5710', '5780']].max(axis=1)
        input_ttl = input_df.sum()

        # Calculate weight
        weight_df = input_df / input_ttl

        # Weighted yield
        weight_yield = yield_df['Yield_5700_5710_5780'] * weight_df

        # Calculate AT yield and output
        AT_Yield = weight_yield.sum()
        AT_out = input_ttl * AT_Yield

        # Calculate ET yield and output
        ET_in = pivot_df['m_In'].get('5600', pd.Series()).sum()
        ET_out = pivot_df['m_Out'].get('5600', pd.Series()).sum()
        ET_yield = ET_out / ET_in if ET_in != 0 else np.nan

        # Adjust yields
        TTL_Yield = 1 - ET_yield * AT_Yield
        ET_yield = 1 - ET_yield
        AT_Yield = 1 - AT_Yield

        # Summarize results
        result = pd.DataFrame({
            'Month': [month_str],  # Add month_str as a new column
            'TTL_Yield': [round(TTL_Yield * 100, 2)],
            'ET_in': [round(ET_in, 0)],
            'ET_out': [round(ET_out, 0)],
            'ET_yield': [round(ET_yield * 100, 2)],
            'AT_in': [round(input_ttl, 0)],
            'AT_out': [round(AT_out, 0)],
            'AT_yield': [round(AT_Yield * 100, 2)]
        })

        return result

    except Exception as e:
        print(f"Error processing data: {e}")
        return None

def get_max_month(db_config):
    """Fetch the maximum month from the database."""
    # 创建 SQLAlchemy 引擎
    engine = create_engine(
        f"mysql+pymysql://{db_config['user']}:{db_config['password']}@{db_config['host']}:{db_config['port']}/{db_config['database']}"
    )
    # 查询数据库
    query = "SELECT MAX(workmt) AS max_month FROM cmsalpha.db_fail_status"
    result = pd.read_sql(query, con=engine)
    return result.iloc[0, 0]

def increment_month(current_month):
    """
    Increment the month, handling year transitions.
    :param current_month: int, the current month in YYYYMM format.
    :return: int, the next month in YYYYMM format.
    """
    year, month = divmod(current_month, 100)
    if month == 12:
        return (year + 1) * 100 + 1  # Next year, January
    else:
        return current_month + 1  # Same year, next month

def Insert_data(result_df, db_config):
    """
    Insert data into the database, and update if the data already exists.
    :param result_df: DataFrame, the data to be inserted.
    :param db_config: dict, database configuration.
    """
    try:
        # 创建 SQLAlchemy 引擎
        engine = create_engine(
            f"mysql+pymysql://{db_config['user']}:{db_config['password']}@{db_config['host']}:{db_config['port']}/{db_config['database']}"
        )

        # Prepare the connection
        with engine.begin() as conn:
            for _, row in result_df.iterrows():
                # Ensure workmt is formatted as a 6-character string (e.g., '202202')
                workmt = str(row['Month']).zfill(6)

                # Replace NaN values with None (for SQL NULL)
                ttl_fail = None if pd.isna(row['TTL_Yield']) else row['TTL_Yield']
                et_in = None if pd.isna(row['ET_in']) else row['ET_in']
                et_out = None if pd.isna(row['ET_out']) else row['ET_out']
                et_yield = None if pd.isna(row['ET_yield']) else row['ET_yield']
                at_in = None if pd.isna(row['AT_in']) else row['AT_in']
                at_out = None if pd.isna(row['AT_out']) else row['AT_out']
                at_yield = None if pd.isna(row['AT_yield']) else row['AT_yield']

                # Build the SQL query string dynamically for each row
                sql = f"""
                INSERT INTO db_fail_status (workmt, ttl_fail, et_in, et_out, et_fail, at_in, at_out, at_fail)
                VALUES ('{workmt}', {ttl_fail}, {et_in}, {et_out}, {et_yield}, {at_in}, {at_out}, {at_yield})
                ON DUPLICATE KEY UPDATE
                    ttl_fail = VALUES(ttl_fail),
                    et_in = VALUES(et_in),
                    et_out = VALUES(et_out),
                    et_fail = VALUES(et_fail),
                    at_in = VALUES(at_in),
                    at_out = VALUES(at_out),
                    at_fail = VALUES(at_fail);
                """

                # Ensure that the SQL string is correct
                # print("Executing SQL query: ", sql)  # Debugging line

                # Use text() to ensure the query is a valid executable object
                conn.execute(text(sql))

        print("Data inserted or updated successfully.")

    except SQLAlchemyError as e:
        print(f"Error inserting or updating data: {e}")


def generate_report_data(db_config, current_date):
    """Generate report data from db_fail_status."""
    try:
        # Connect to the database
        engine = create_engine(
            f"mysql+pymysql://{db_config['user']}:{db_config['password']}@{db_config['host']}:{db_config['port']}/{db_config['database']}"
        )

        # Calculate the years and months for the report
        current_year = current_date.year
        current_month = current_date.month

        # Generate a list of the last 3 years
        years = [current_year - 1, current_year - 2, current_year - 3]

        # Get the last 12 months
        months = [
            (current_year, current_month - i) if current_month - i > 0 else (current_year - 1, 12 + (current_month - i))
            for i in range(1, 13)]

        # Generate the query to get the data
        # Query for the last 3 years
        year_query = """
                    SELECT 
                    YEAR(STR_TO_DATE(workmt, '%Y%m')) AS year,
                    (1 - SUM(et_out) / NULLIF(SUM(et_in), 0) * SUM(at_out) / NULLIF(SUM(at_in), 0)) * 100 AS ttl_fail,
                    SUM(et_in) AS et_in,
                    SUM(et_out) AS et_out,
                    (1 - SUM(at_out) / NULLIF(SUM(at_in), 0)) * 100 AS et_fail,
                    SUM(at_in) AS at_in,
                    SUM(at_out) AS at_out,
                    (1 - SUM(at_out) / NULLIF(SUM(at_in), 0)) * 100 AS at_fail
                    FROM db_fail_status
                    WHERE YEAR(STR_TO_DATE(workmt, '%Y%m')) IN :years
                    GROUP BY YEAR(STR_TO_DATE(workmt, '%Y%m'))
        """

        # Execute the query to get the last 3 years of data
        with engine.connect() as conn:
            result_years = conn.execute(text(year_query), {"years": tuple(years)}).fetchall()

        # Convert the result to a DataFrame for easy manipulation
        df_years = pd.DataFrame(result_years,
                                columns=['year', 'ttl_fail', 'et_in', 'et_out', 'et_fail', 'at_in', 'at_out',
                                         'at_fail'])

        # Get the data for the last 12 months
        month_query = """
        SELECT 
            workmt,
            ttl_fail,
            et_in,
            et_out,
            et_fail,
            at_in,
            at_out,
            at_fail
        FROM db_fail_status
        WHERE workmt BETWEEN :start_date AND :end_date
        """

        # Calculate the start and end date for the last 12 months
        start_date = f"{months[-1][0]}-{months[-1][1]:02d}-01"
        end_date = f"{current_date.year}-{current_date.month - 1 if current_date.month > 1 else 12:02d}-01"

        # Execute the query to get the last 12 months data
        with engine.connect() as conn:
            result_months = conn.execute(text(month_query), {"start_date": start_date, "end_date": end_date}).fetchall()

        # Convert the result to a DataFrame for easy manipulation
        df_months = pd.DataFrame(result_months,
                                 columns=['workmt', 'ttl_fail', 'et_in', 'et_out', 'et_fail', 'at_in', 'at_out',
                                          'at_fail'])

        # Combine the year and month data
        # Combine the 3 years data with the 12 months data (total 15 rows)
        report_data = pd.concat([df_years, df_months], ignore_index=True)

        # Return the report data as a dictionary
        report_dict = report_data.to_dict(orient='records')

        # Convert report data into a DataFrame
        df = pd.DataFrame(report_data)

        # Combine 'year' and 'workmt' into a single column 'period'
        df['period'] = df['workmt'].combine_first(df['year'].apply(lambda x: f"{int(x):04}" if pd.notna(x) else None))

        # Drop the original 'year' and 'workmt' columns
        df = df.drop(columns=['year', 'workmt'])

        # Reorder columns to place 'period' at the front
        cols = ['period'] + [col for col in df.columns if col != 'period']
        df = df[cols]

        return df

    except Exception as e:
        print(f"Error generating report data: {e}")
        return None


def update_chart_in_ppt(ppt_file, chart_slide_index, chart_index, categories, values):
    """
    Update chart data in a specific slide of a PowerPoint presentation.

    :param ppt_file: str, path to the PowerPoint file.
    :param chart_slide_index: int, index of the slide containing the chart.
    :param chart_index: int, index of the chart within the slide.
    :param categories: list, categories for the chart (e.g., months or years).
    :param values: list, values corresponding to the categories.
    """
    # Load the presentation
    prs = Presentation(ppt_file)

    # Access the slide containing the chart
    slide = prs.slides[chart_slide_index]

    # Get the chart
    chart = slide.shapes[chart_index].chart

    # Update the chart with new data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Series 1', values)

    # Replace the existing chart
    x, y, cx, cy = chart.shape.left, chart.shape.top, chart.shape.width, chart.shape.height
    slide.shapes._spTree.remove(chart._element)  # Remove the old chart
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

    # Save changes
    prs.save("updated_" + ppt_file)
    print(f"Updated chart data saved to updated_{ppt_file}")


def update_table_in_ppt(ppt_file, data, table_slide_index=0, table_index=2):
    """
    Update table data in a specific slide of a PowerPoint presentation (transposed format).
    :param ppt_file: str, path to the PowerPoint file.
    :param data: list of lists, data to be written into the table.
                 Each sublist represents a column, as the data is transposed.
    :param table_slide_index: int, index of the slide containing the table.
    :param table_index: int, index of the table within the slide.
    """
    # Load the presentation
    prs = Presentation(ppt_file)

    # Access the slide containing the table
    slide = prs.slides[table_slide_index]

    # Get the table shape by index
    table_shape = slide.shapes[table_index]
    if table_shape.shape_type != MSO_SHAPE_TYPE.TABLE:
        raise ValueError(f"The shape at index {table_index} is not a table.")

    # Access the table
    table = table_shape.table

    data = pd.DataFrame(data)
    # Transpose the DataFrame to align rows and columns for the table
    transposed_data = data.transpose()

    # Iterate over rows in the transposed data
    for row_idx, (index, row) in enumerate(transposed_data.iterrows()):
        for col_idx, value in enumerate(row):
            # Adjust the target cell positions
            table_row_idx = row_idx   # Data rows start from the second row in the table
            table_col_idx = col_idx + 1  # Data columns start from the third column in the table

            # Format the value
            if pd.isna(value):
                value = ""
            elif "Fail" in transposed_data.iloc[0, col_idx]:  # Format fail rates
                value = f"{value * 100:.2f}%"
            elif isinstance(value, (int, float)):  # Format in and out values
                value = f"{int(value):,}"

            # Write to table cell
            print(f"Writing {value} to cell ({table_row_idx}, {table_col_idx})")
            cell = table.cell(table_row_idx, table_col_idx)
            cell.text = str(value)

            # Apply formatting
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)  # Adjust font size
            paragraph.alignment = PP_ALIGN.CENTER  # Center align

            if "Fail" in transposed_data.iloc[0, col_idx]:  # Bold fail rates
                paragraph.font.bold = True

    # Save the updated presentation
    updated_file = ppt_file
    prs.save(updated_file)
    print(f"Updated table saved to {updated_file}")


def write_to_ppt(result, reportDir):
    # Convert result data to table format
    table_data = [['Period', 'TTL Fail', 'ET In', 'ET Out', 'ET Fail', 'AT In', 'AT Out', 'AT Fail']]
    for _, row in result.iterrows():
        table_data.append([
            row['period'] if 'period' in row else '',
            row['ttl_fail'] if 'ttl_fail' in row else 0,
            row['et_in'] if 'et_in' in row else 0,
            row['et_out'] if 'et_out' in row else 0,
            row['et_fail'] if 'et_fail' in row else 0,
            row['at_in'] if 'at_in' in row else 0,
            row['at_out'] if 'at_out' in row else 0,
            row['at_fail'] if 'at_fail' in row else 0,
        ])

    # Update table in slide 0, table index 0
    update_table_in_ppt(reportDir, table_data, table_slide_index=0, table_index=2)

    # Update chart in slide 1, chart index 1 (e.g., only last 12 months)
    # categories = [record['workmt'] for record in result[-12:]]
    # values = [record['ttl_fail'] for record in result[-12:]]
    # update_chart_in_ppt(reportDir, chart_slide_index=1, chart_index=0, categories=categories, values=values)

def list_shapes_in_slide(ppt_file, slide_index):
    """List all shapes and their types in a slide."""
    prs = Presentation(ppt_file)
    slide = prs.slides[slide_index]
    for idx, shape in enumerate(slide.shapes):
        print(f"Index: {idx}, Type: {shape.shape_type}, Name: {shape.name}")

def main(mode):
    # 确定数据库配置
    if mode == 'test':
        db_config = {
        'host': 'localhost',
        'user': 'remoteuser',
        'password': 'password',
        'database': 'cmsalpha',
        'charset': 'utf8mb4',
        'port': 3306,
        }
        reportDir = r'D:\Sync\业务报告\1on1\2025\1月\不良Status.pptx'
    else:
        db_config = {
            'host': '172.27.154.57',
            'user': 'remoteuser',
            'password': 'password',
            'database': 'cmsalpha',
            'charset': 'utf8mb4',
            'port': 3306,
        }
        reportDir = r'\\172.27.7.199\Mod_TestE\20. Fail Status\不良Status.pptx'
    # 确定数据库最大月份
    max_month = get_max_month(db_config)
    # print(f"Max month in the database: {max_month}")

    # 获取当前系统月份的前一个月
    current_date = datetime.now()
    end_month = (current_date.replace(day=1) - timedelta(days=1)).strftime("%Y%m")
    max_month = '202411'
    # 比较 end_month 和 max_month，并逐月执行
    if end_month > max_month:
        # 从 max_month 的下一个月份开始
        current_month = increment_month(int(max_month))

        while current_month <= int(end_month):
            month_str = str(current_month)
            # print(month_str)
            result_df = process_data_from_db(db_config, month_str)
            if result_df is not None:
                print(f"Results for {month_str}:")
                print(result_df)
                # 将结果写入数据库
                Insert_data(result_df, db_config)
            # 增加月份
            current_month = increment_month(int(current_month))
            # 生成报表数据
            result = generate_report_data(db_config, current_date)
            print("Report data:", result)
            # 写入PPT
            list_shapes_in_slide(reportDir, 0)
            write_to_ppt(result, reportDir)
    else:
        print("No new months to process.")


# Main function to execute processing
if __name__ == "__main__":
    main('test')