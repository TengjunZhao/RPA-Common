import pandas as pd
import numpy as np
import pymysql
from datetime import datetime, timedelta
import calendar
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from decimal import Decimal
from pptx.util import Pt
import shutil
import os



def get_first_and_last_day_of_month(month_str):
    """
    Get the first and last day of the month given a string in YYYYMM format.
    :param month_str: str, the month in YYYYMM format.
    :return: tuple, (first_day, last_day) in YYYY-MM-DD format.
    """
    year = int(month_str[:4])
    month = int(month_str[4:])

    # First day of the month
    first_day = f"{year}{month:02d}01"

    # Last day of the month
    last_day = f"{year}{month:02d}{calendar.monthrange(year, month)[1]}"

    return first_day, last_day

def fetch_data_from_db(connection, query):
    """Fetch data from the database using the provided query."""
    return pd.read_sql(query, connection)

def process_data_from_db(db_config, month_str):
    """Process data from the database and compute the required metrics."""
    first_day, last_day = get_first_and_last_day_of_month(month_str)
    try:
        # Connect to the database using pymysql
        connection = pymysql.connect(
            host=db_config['host'],
            user=db_config['user'],
            password=db_config['password'],
            database=db_config['database'],
            port=db_config['port'],
            charset=db_config['charset']
        )

        # Query to fetch required data
        query = f"""
        SELECT 
            dy.modtype,
            dd.Product_density,
            dy.oper_old, 
            SUM(dy.in_qty) AS m_In, 
            SUM(dy.out_qty) AS m_Out
        FROM 
            cmsalpha.db_yielddetail dy
        JOIN modulemte.db_deviceinfo dd on dy.device = dd.Device
        WHERE 
            workdt BETWEEN '{first_day}' AND '{last_day}' AND
            dy.oper_old in ('5600', '5700', '5710', '5780')
        GROUP BY 
            dy.modtype, dd.Product_density, oper_old
        """
        # Fetch data
        data = fetch_data_from_db(connection, query)

        # Pivot table后处理列名
        pivot_df = data.pivot_table(index=['modtype', 'Product_density'], columns='oper_old', values=['m_In', 'm_Out'],
                                    aggfunc='sum')

        # 扁平化多级列名
        pivot_df.columns = [f"{metric}_{oper}" for metric, oper in pivot_df.columns]
        pivot_df = pivot_df.fillna(0)

        # 计算5600指标（ET测试）
        pivot_df['IN_5600'] = pivot_df['m_In_5600']
        pivot_df['OUT_5600'] = pivot_df['m_Out_5600']
        pivot_df['Fail rate_5600'] = round((1 - (pivot_df['OUT_5600'] / pivot_df['IN_5600'])) * 100, 2)
        pivot_df.loc[pivot_df['IN_5600'] == 0, 'Fail rate_5600'] = np.nan  # 处理除零情况

        # 计算5710-5700-5780指标（AT测试）
        pivot_df['IN_5710_5700_5780'] = pivot_df['m_In_5700'] + pivot_df['m_In_5710'] + pivot_df['m_In_5780']
        pivot_df['OUT_5710_5700_5780'] = pivot_df['m_Out_5700'] + pivot_df['m_Out_5710'] + pivot_df['m_Out_5780']
        pivot_df['Fail rate_5710_5700_5780'] = round((1 - (pivot_df['OUT_5710_5700_5780'] /
                                                           pivot_df['IN_5710_5700_5780'])) * 100, 2)
        pivot_df.loc[pivot_df['IN_5710_5700_5780'] == 0, 'Fail rate_5710_5700_5780'] = np.nan

        # Summarize results
        result = pivot_df[[
            'IN_5600', 'OUT_5600', 'Fail rate_5600',
            'IN_5710_5700_5780', 'OUT_5710_5700_5780', 'Fail rate_5710_5700_5780'
        ]].reset_index()

        # 添加workmt字段并重命名列
        result['workmt'] = month_str  # 新增月份字段
        result = result.rename(columns={
            'Product_density': 'density',
            'IN_5600': 'et_in',
            'OUT_5600': 'et_out',
            'Fail rate_5600': 'et_fail',
            'IN_5710_5700_5780': 'at_in',
            'OUT_5710_5700_5780': 'at_out',
            'Fail rate_5710_5700_5780': 'at_fail'
        })

        # 计算总不良率（ET和AT不良率的乘积）
        result[['et_in', 'et_out', 'at_in', 'at_out', 'et_fail', 'at_fail']] = (
            result[['et_in', 'et_out', 'at_in', 'at_out', 'et_fail', 'at_fail']].fillna(0))
        result['ttl_fail'] = (1 - (1 - result['et_fail']) * (1 - result['at_fail']))
        result['ttl_fail'] = result['ttl_fail'].round(2)

        # 选择最终输出列
        result = result[[
            'workmt', 'modtype', 'density',
            'ttl_fail',
            'et_in', 'et_out', 'et_fail',
            'at_in', 'at_out', 'at_fail'
        ]]
        return result

    except Exception as e:
        print(f"Error processing data: {e}")
        return None

    finally:
        # Close the connection after use
        if connection:
            connection.close()

def get_max_month(db_config):
    """Fetch the maximum month from the database."""
    connection = None  # Ensure connection is initialized
    try:
        # Connect to the database using pymysql
        connection = pymysql.connect(
            host=db_config['host'],
            user=db_config['user'],
            password=db_config['password'],
            database=db_config['database'],
            port=db_config['port'],
            charset=db_config['charset']
        )

        query = "SELECT MAX(workmt) AS max_month FROM cmsalpha.db_fail_status_density"
        result = pd.read_sql(query, con=connection)
        return result.iloc[0, 0]

    except Exception as e:
        print(f"Error fetching max month: {e}")
        return None
    finally:
        if connection:
            connection.close()

def increment_month(current_month):
    """
    Increment the month, handling year transitions.
    :param current_month: int, the current month in YYYYMM format.
    :return: int, the next month in YYYYMM format.
    """
    year, month = divmod(current_month, 100)
    if month == 12:
        return (year + 1) * 100 + 1  # Next year, January
    else:
        return current_month + 1  # Same year, next month

def Insert_data(result_df, db_config):
    """Insert data into the database, and update if the data already exists."""
    try:
        # Connect to the database using pymysql
        connection = pymysql.connect(
            host=db_config['host'],
            user=db_config['user'],
            password=db_config['password'],
            database=db_config['database'],
            port=db_config['port'],
            charset=db_config['charset']
        )

        with connection.cursor() as cursor:
            for _, row in result_df.iterrows():
                # Ensure workmt is formatted as a 6-character string (e.g., '202202')
                workmt = str(row['workmt']).zfill(6)

                # Replace NaN values with None (for SQL NULL)
                modtype = None if pd.isna(row['modtype']) else row['modtype']
                density = None if pd.isna(row['density']) else row['density']
                ttl_fail = None if pd.isna(row['ttl_fail']) else row['ttl_fail']
                et_in = None if pd.isna(row['et_in']) else row['et_in']
                et_out = None if pd.isna(row['et_out']) else row['et_out']
                et_yield = None if pd.isna(row['et_fail']) else row['et_fail']
                at_in = None if pd.isna(row['at_in']) else row['at_in']
                at_out = None if pd.isna(row['at_out']) else row['at_out']
                at_yield = None if pd.isna(row['at_fail']) else row['at_fail']

                # Build the SQL query string dynamically for each row
                sql = f"""
                INSERT INTO db_fail_status_density (workmt,modtype, modDensity, ttl_fail, et_in, et_out, et_fail, at_in, at_out, at_fail)
                VALUES ('{workmt}','{modtype}', '{density}', {ttl_fail}, {et_in}, {et_out}, {et_yield}, {at_in}, {at_out}, {at_yield})
                ON DUPLICATE KEY UPDATE
                    ttl_fail = VALUES(ttl_fail),
                    et_in = VALUES(et_in),
                    et_out = VALUES(et_out),
                    et_fail = VALUES(et_fail),
                    at_in = VALUES(at_in),
                    at_out = VALUES(at_out),
                    at_fail = VALUES(at_fail);
                """
                cursor.execute(sql)

            connection.commit()
            print("Data inserted or updated successfully.")

    except pymysql.MySQLError as e:
        print(f"Error inserting or updating data: {e}")

    finally:
        if connection:
            connection.close()

def generate_report_failStatus(db_config, current_date, product):
    """Generate report data from db_fail_status using pymysql."""
    try:
        # 连接到数据库
        connection = pymysql.connect(
            host=db_config['host'],
            user=db_config['user'],
            password=db_config['password'],
            database=db_config['database'],
            port=db_config['port']
        )

        # 计算报告的年份和月份
        current_year = current_date.year
        current_month = current_date.month

        # 获取过去3年的年份
        years = [current_year, current_year - 1, current_year - 2]

        # 获取过去12个月的月份
        months = [
            (current_year, current_month - i) if current_month - i > 0 else (current_year - 1, 12 + (current_month - i))
            for i in range(1, 13)
        ]

        # 年份数据的查询
        # 将查询中的 IN 部分改成多个占位符
        year_query = """
                SELECT 
                    YEAR(STR_TO_DATE(workmt, '%%Y%%m')) AS year,
                    modtype, 
                    modDensity,
                    (1 - SUM(et_out) / NULLIF(SUM(et_in), 0) * SUM(at_out) / NULLIF(SUM(at_in), 0)) * 100 AS ttl_fail,
                    SUM(et_in) AS et_in,
                    SUM(et_out) AS et_out,
                    (1 - SUM(et_out) / NULLIF(SUM(et_in), 0)) * 100 AS et_fail,
                    SUM(at_in) AS at_in,
                    SUM(at_out) AS at_out,
                    (1 - SUM(at_out) / NULLIF(SUM(at_in), 0)) * 100 AS at_fail
                FROM db_fail_status_density
                WHERE YEAR(STR_TO_DATE(workmt, '%%Y%%m')) IN ({}) 
                    AND modtype = %s 
                    AND modDensity = %s
                GROUP BY YEAR(STR_TO_DATE(workmt, '%%Y%%m')), modtype, modDensity
            """.format(
            ','.join(['%s'] * len(years))
        )

        # 打印调试信息
        print("Years to query:", tuple(years))  # 打印 years 参数
        # 执行查询时传入合并参数
        query_params = tuple(years) + (product[0], product[1])
        # 执行查询
        with connection.cursor() as cursor:
            cursor.execute(year_query, query_params)  # 传入实际的年份元组
            result_years = cursor.fetchall()
            print("debug2", result_years)

        # 将结果转换为 DataFrame
        df_years = pd.DataFrame(result_years, columns=[
            'year', 'modtype', 'density', 'ttl_fail',
            'et_in', 'et_out', 'et_fail',
            'at_in', 'at_out', 'at_fail'
        ])

        # 获取过去12个月的数据
        month_query = """
                SELECT 
                    workmt,
                    modtype,
                    modDensity,
                    ttl_fail,
                    et_in,
                    et_out,
                    et_fail,
                    at_in,
                    at_out,
                    at_fail
                FROM db_fail_status_density
                WHERE workmt BETWEEN %s AND %s AND
                    modtype = %s AND  # 改为单个占位符
                    modDensity = %s   # 改为单个占位符
                """

        # 计算开始和结束日期
        start_date = f"{months[-1][0]}{months[-1][1]:02d}"
        end_date = f"{current_date.year}{current_date.month - 1 if current_date.month > 1 else 12:02d}"

        # 执行查询并获取结果
        with connection.cursor() as cursor:
            query_params = (start_date, end_date) + (product[0], product[1])
            cursor.execute(month_query, query_params)
            result_months = cursor.fetchall()

        # 将结果转换为 DataFrame
        df_months = pd.DataFrame(result_months, columns=[
            'workmt','modtype','modDensity', 'ttl_fail', 'et_in', 'et_out', 'et_fail', 'at_in', 'at_out', 'at_fail'
        ])

        # 合并年度数据和月度数据
        report_data = pd.concat([df_years, df_months], ignore_index=True)

        # 处理period列（合并year和workmt）
        report_data['period'] = report_data['workmt'].combine_first(
            report_data['year'].apply(lambda x: f"{int(x):04}" if pd.notna(x) else None)
        )
        report_data['density'] = report_data['density'].combine_first(report_data['modDensity'])

        # 先对年度数据排序（升序）
        df_years_sorted = df_years.sort_values(by='year', ascending=True)

        # 对月度数据排序（按workmt升序）
        df_months_sorted = df_months.sort_values(by='workmt', ascending=True)

        # 重新合并排序后的数据
        sorted_data = pd.concat([df_years_sorted, df_months_sorted], ignore_index=True)

        # 处理period列
        sorted_data['period'] = sorted_data['workmt'].combine_first(
            sorted_data['year'].apply(lambda x: f"{int(x):04}" if pd.notna(x) else None)
        )
        sorted_data['density'] = sorted_data['density'].combine_first(sorted_data['modDensity'])

        # 删除不需要的列
        sorted_data = sorted_data.drop(columns=['year', 'workmt', 'modDensity'])

        # 重新排列列顺序
        cols = ['period'] + [col for col in sorted_data.columns if col != 'period']
        sorted_data = sorted_data[cols]

        # 创建完整时间段索引（保持原始顺序）
        all_periods = [f"{y}" for y in sorted(years)] + [f"{m[0]}{m[1]:02}" for m in
                                                         sorted(months, key=lambda x: (x[0], x[1]))]
        multi_index = pd.MultiIndex.from_product(
            [all_periods, sorted_data['modtype'].unique(), sorted_data['density'].unique()],
            names=['period', 'modtype', 'density']
        )

        # 重新索引并填充缺失值
        final_df = sorted_data.set_index(['period', 'modtype', 'density']).reindex(multi_index).fillna(0).reset_index()

        # 数值列保留两位小数
        for col in ['ttl_fail', 'et_fail', 'at_fail']:
            final_df[col] = final_df[col].apply(lambda x: round(x, 2))

        return final_df

    except Exception as e:
        print(f"Error generating report data: {e}")
        return None
    finally:
        # 确保连接被关闭
        if connection:
            connection.close()

# 获取产品的target(以当月前三个月均值为基准)
def generate_target(db_config, df, product):
    try:
        # 连接到数据库
        connection = pymysql.connect(
            host=db_config['host'],
            user=db_config['user'],
            password=db_config['password'],
            database=db_config['database'],
            port=db_config['port']
        )
        # df增加一列targeet
        df['target'] = 0
        # 逐行处理df数据
        for index, row in df.iterrows():
            period = row['period']
            modtype = row['modtype']
            density = row['density']
            # 如果period是年份
            if len(period) == 4:
                # 查询period所在年份，modtype，density，的ttl_fail的平均值
                query = f"""
                    SELECT ROUND(AVG(ttl_fail),2) AS avg_ttl_fail
                    FROM db_fail_status_density
                    WHERE SUBSTRING(workmt,1,4) = '{period}'
                        AND modtype = '{modtype}'
                        AND modDensity = '{density}'
                """
                with connection.cursor() as cursor:
                    cursor.execute(query)
                    result = cursor.fetchone()
                    if result:
                        avg_ttl_fail = result[0]
                        df.at[index, 'target'] = avg_ttl_fail
                    else:
                        df.at[index, 'target'] = 0

            else:
                # 查询period所在月份的前三个月，modtype，density的ttl_fail的平均值
                query = """
                    SELECT ROUND(AVG(ttl_fail), 2) AS avg_ttl_fail
                    FROM db_fail_status_density
                    WHERE workmt BETWEEN %s AND %s
                        AND modtype = %s
                        AND modDensity = %s
                """

                # 计算开始和结束月份
                year = int(period[0:4])
                month = int(period[4:6])
                # 计算结束月份（前一个月）
                if month > 1:
                    end_month = month - 1
                    end_year = year
                else:
                    end_month = 12
                    end_year = year - 1
                # 计算前三个月的起始月份
                if month > 3:
                    start_month = month - 3
                    start_year = year
                else:
                    start_month = 12 + (month - 3)
                    start_year = year - 1

                # 格式化日期
                start_date = f"{start_year}{start_month:02d}"
                end_date = f"{end_year}{end_month:02d}"

                # 执行查询
                with connection.cursor() as cursor:
                    cursor.execute(query, (start_date, end_date, modtype, density))
                    result = cursor.fetchone()
                    avg_ttl_fail = result[0] if result else 0.0
                    df.at[index, 'target'] = avg_ttl_fail
        return df

    except Exception as e:
        print(f"Error generating report data: {e}")
        return None
    finally:
        # 确保连接被关闭
        if connection:
            connection.close()

# 获取产品的target(以当月上季度均值为基准)
def generate_target_season(db_config, df, product):
    try:
        # 连接到数据库
        connection = pymysql.connect(
            host=db_config['host'],
            user=db_config['user'],
            password=db_config['password'],
            database=db_config['database'],
            port=db_config['port']
        )
        # df增加一列target
        df['target'] = 0
        # 逐行处理df数据
        for index, row in df.iterrows():
            period = row['period']
            modtype = row['modtype']
            density = row['density']
            # 如果period是年份
            if len(period) == 4:
                # 查询period所在年份，modtype，density，的ttl_fail的平均值
                query = f"""
                    SELECT ROUND(AVG(ttl_fail),2) AS avg_ttl_fail
                    FROM db_fail_status_density
                    WHERE SUBSTRING(workmt,1,4) = '{period}'
                        AND modtype = '{modtype}'
                        AND modDensity = '{density}'
                """
                with connection.cursor() as cursor:
                    cursor.execute(query)
                    result = cursor.fetchone()
                    if result:
                        avg_ttl_fail = result[0]
                        df.at[index, 'target'] = avg_ttl_fail
                    else:
                        df.at[index, 'target'] = 0

            else:
                # 查询period所在月份的上一季度的平均值
                query = """
                    SELECT ROUND(AVG(ttl_fail), 2) AS avg_ttl_fail
                    FROM db_fail_status_density
                    WHERE workmt BETWEEN %s AND %s
                        AND modtype = %s
                        AND modDensity = %s
                """

                # 计算季度信息
                year = int(period[0:4])
                month = int(period[4:6])

                # 确定当前季度和上一季度
                current_quarter = (month - 1) // 3 + 1
                if current_quarter == 1:
                    prev_quarter = 4
                    prev_year = year - 1
                else:
                    prev_quarter = current_quarter - 1
                    prev_year = year

                # 计算上一季度的开始和结束月份
                start_month = (prev_quarter - 1) * 3 + 1
                end_month = prev_quarter * 3

                # 格式化日期
                start_date = f"{prev_year}{start_month:02d}"
                end_date = f"{prev_year}{end_month:02d}"

                # 执行查询
                with connection.cursor() as cursor:
                    cursor.execute(query, (start_date, end_date, modtype, density))
                    result = cursor.fetchone()
                    avg_ttl_fail = result[0] if result else 0.0
                    df.at[index, 'target'] = avg_ttl_fail
        return df

    except Exception as e:
        print(f"Error generating report data: {e}")
        return None
    finally:
        # 确保连接被关闭
        if connection:
            connection.close()

# 获取产品的target（以当年为目标）
def generate_target_year(db_config, df, product):
    try:
        # 连接到数据库
        connection = pymysql.connect(
            host=db_config['host'],
            user=db_config['user'],
            password=db_config['password'],
            database=db_config['database'],
            port=db_config['port']
        )
        # df增加一列target
        df['target'] = 0
        # 逐行处理df数据
        for index, row in df.iterrows():
            period = row['period']
            modtype = row['modtype']
            density = row['density']
            # 如果period是年份
            if len(period) == 4:
                # 查询period所在年份，modtype，density，的ttl_fail的平均值
                query = f"""
                    SELECT ROUND(AVG(ttl_fail),2) AS avg_ttl_fail
                    FROM db_fail_status_density
                    WHERE SUBSTRING(workmt,1,4) = '{period}'
                        AND modtype = '{modtype}'
                        AND modDensity = '{density}'
                """
                with connection.cursor() as cursor:
                    cursor.execute(query)
                    result = cursor.fetchone()
                    if result:
                        avg_ttl_fail = result[0]
                        df.at[index, 'target'] = avg_ttl_fail
                    else:
                        df.at[index, 'target'] = 0
                latest = df.at[index, 'target']
                print(latest)
            else:
                df.at[index, 'target'] = latest
        return df

    except Exception as e:
        print(f"Error generating report data: {e}")
        return None
    finally:
        # 确保连接被关闭
        if connection:
            connection.close()

def update_chart_in_ppt(ppt_file, result, chart_slide_index, chart_index):
    """
    Update chart data in a specific slide of a PowerPoint presentation.
    :param ppt_file: str, path to the PowerPoint file.
    :param result: pandas DataFrame, result data for the chart.
    :param chart_slide_index: int, index of the slide containing the chart.
    :param chart_index: int, index of the chart within the slide.
    """
    # Load the presentation
    prs = Presentation(ppt_file)

    # Access the slide containing the chart
    slide = prs.slides[chart_slide_index]

    # Get the chart shape by index
    chart_shape = slide.shapes[chart_index]
    if chart_shape.shape_type != MSO_SHAPE_TYPE.CHART:
        raise ValueError(f"The shape at index {chart_index} is not a chart.")

    # Access the chart
    chart = chart_shape.chart

    # Prepare data for the chart (e.g., TTL Fail for the last 12 months)
    # Extracting the last 12 periods and TTL Fail values
    months = result['period'].unique().tolist()
    formatted_months = [
        f"{p[0:4]}" if len(p) == 4 else f"{p[2:4]}년 {p[4:]}월"
        for p in sorted(months, key=lambda x: (len(x), x))
    ]
    # 创建包含4个系列的新图表数据
    chart_data = CategoryChartData()
    chart_data.categories = formatted_months

    # ET Fail Rate（柱状图）
    chart_data.add_series('ET Fail Rate',
                          [round(val, 2) for val in result['et_fail'].fillna(0).tolist()])

    # AT Fail Rate（柱状图）
    chart_data.add_series('AT Fail Rate',
                          [round(val, 2) for val in result['at_fail'].fillna(0).tolist()])

    # Target（折线图）
    chart_data.add_series('Target',
                          [round(val, 2) if pd.notnull(val) else 0 for val in result['target'].tolist()])

    # Replace the existing data with the new chart_data
    chart.replace_data(chart_data)

    # Save the updated presentation
    updated_file = ppt_file
    prs.save(updated_file)
    print(f"Updated chart data saved to {updated_file}")


def update_table_in_ppt(ppt_file, data, table_slide_index=0, table_index=2):
    """
    Update table data in a specific slide of a PowerPoint presentation (transposed format).
    :param ppt_file: str, path to the PowerPoint file.
    :param data: list of lists, data to be written into the table.
                 Each sublist represents a column, as the data is transposed.
    :param table_slide_index: int, index of the slide containing the table.
    :param table_index: int, index of the table within the slide.
    """
    # Load the presentation
    prs = Presentation(ppt_file)

    # Access the slide containing the table
    slide = prs.slides[table_slide_index]

    # Get the table shape by index
    table_shape = slide.shapes[table_index]
    if table_shape.shape_type != MSO_SHAPE_TYPE.TABLE:
        raise ValueError(f"The shape at index {table_index} is not a table.")

    # Access the table
    table = table_shape.table

    # 创建密度分类的层级结构
    density_groups = data.groupby('density')

    # 动态生成时间段（从数据中提取年份和月份）
    all_periods = data['period'].unique().astype(str)
    # 分离年份和月份数据
    # 分离并排序原始数据（保持原始格式）
    years = sorted([p for p in all_periods if len(p) == 4])
    month_periods = sorted([p for p in all_periods if len(p) == 6], key=lambda x: x)
    # 合并原始顺序（此时已排序）
    raw_periods = years + month_periods
    # 转换格式（仅在最后阶段处理一次）
    formatted_periods = [
        p if len(p) == 4 else f"{p[2:4]}-{calendar.month_abbr[int(p[4:])]}"
        for p in raw_periods
    ]
    # 生成时间段表头（第一行）
    header_row = ["Density",'Period'] + formatted_periods
    table_data = [header_row]
    density_order = ['4GB', '8GB', '16GB', '32GB', '64GB', '128GB']

    for density in density_order:  # 按指定顺序处理密度
        if density not in density_groups.groups:
            continue

        # 获取并排序当前密度数据
        group = pd.DataFrame(density_groups.get_group(density)).copy()
        group['sort_order'] = group['period'].apply(lambda x: raw_periods.index(x))
        group.sort_values('sort_order', inplace=True)
        if group.empty:
            print(f"警告：{density}分组数据为空")
            continue
        # # 按时间段顺序重新索引数据
        # ordered_group = group.set_index('period').reindex(periods).fillna(0)

        # 生成密度区块
        table_data += [
            [density, "Rate(%)"] + group['ttl_fail'].map("{:.2f}%".format).tolist(),
            ["", "ET IN"] + group['et_in'].map("{:,.0f}".format).tolist(),
            ["", "ET OUT"] + group['et_out'].map("{:,.0f}".format).tolist(),
            ["", "ET Rate(%)"] + group['et_fail'].map("{:.2f}%".format).tolist(),
            ["", "AT IN"] + group['at_in'].map("{:,.0f}".format).tolist(),
            ["", "AT OUT"] + group['at_out'].map("{:,.0f}".format).tolist(),
            ["", "AT Rate(%)"] + group['at_fail'].map("{:.2f}%".format).tolist()
        ]
    max_cols = len(header_row)
    for row_idx, row in enumerate(table_data):
        # 自动补全缺失列
        padded_row = list(row) + [''] * (max_cols - len(row))
        # 判断是否为Rate行（通过第二列特征）
        is_rate_row = any(keyword in row[1] for keyword in ['Rate', 'SUM Rate']) if len(row) > 1 else False
        for col_idx, value in enumerate(padded_row[:max_cols]):
            if col_idx < 2:
                continue
            try:
                cell = table.cell(row_idx, col_idx + 1)
                cell.text = str(value)
                # 统一格式设置
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(10)
                paragraph.alignment = PP_ALIGN.CENTER
                # 整行加粗逻辑
                if is_rate_row:
                    paragraph.font.bold = True
            except IndexError:
                print(f"表格尺寸不足，需要至少 {row_idx + 1} 行 {col_idx + 1} 列")

    prs.save(ppt_file)
    print(f"Updated table saved to {ppt_file}")

def write_to_ppt(result, reportDir, product, sheet):
    if result is None:
        print("No report data available to write to PPT.")
        return
    # 根据product在products列表中的位置确定chart_index
    products = [('SD', '4GB'), ('SD', '8GB'), ('SD', '16GB'),
                ('RD', '16GB'), ('RD', '32GB'), ('RD', '64GB')]
    chart_index = products.index(product)
    # Update table in slide 0, table index 0
    # update_table_in_ppt(reportDir, filtered_data, table_slide_index=sheet, table_index=2)

    # Update chart in slide 1, chart index 1 (e.g., only last 12 months)
    update_chart_in_ppt(reportDir, result, chart_slide_index=sheet, chart_index=chart_index)


def list_shapes_in_slide(ppt_file, slide_index):
    """List all shapes and their types in a slide."""
    prs = Presentation(ppt_file)
    slide = prs.slides[slide_index]
    for idx, shape in enumerate(slide.shapes):
        print(f"Index: {idx}, Type: {shape.shape_type}, Name: {shape.name}")

def copy_and_save_report(src_path, target_dir, report_month):
    threshold_date = datetime.strptime(report_month, '%Y%m')  # Convert report month to datetime
    threshold_month = (threshold_date.replace(year=threshold_date.year - 1)).strftime('%Y%m')
    for folder_name in os.listdir(target_dir):
        folder_path = os.path.join(target_dir, folder_name)
        if os.path.isdir(folder_path):  # Check if it's a directory
            try:
                # Check if folder name is a valid YYYYMM format and compare
                if folder_name.isdigit() and len(folder_name) == 6:
                    if folder_name < threshold_month:
                        shutil.rmtree(folder_path)  # Delete the folder and its contents
                        print(f"Deleted old report folder: {folder_path}")
            except Exception as e:
                print(f"Error deleting folder {folder_path}: {e}")
    # Create the new target directory
    new_dir = os.path.join(target_dir, report_month)  # Combine base directory with YYYYMM
    os.makedirs(new_dir, exist_ok=True)  # Create the directory if it doesn't exist

    # Construct the new file path
    new_file_path = os.path.join(new_dir, os.path.basename(src_path))

    # Copy the file to the new location (overwrite if exists)
    shutil.copy2(src_path, new_file_path)

    print(f"Report successfully copied to: {new_file_path}")


def main(mode):
    pcType = ['SD', 'UD']
    pptPCType = ['SD']
    svType = ['RD', 'LD']
    pptSVType = ['RD']
    pcDensity = ['4GB', '8GB', '16GB']
    svDensity = ['16GB', '32GB', '64GB']
    products = [('SD','4GB'), ('SD','8GB'),('SD','16GB'),
               ('RD','16GB'), ('RD','32GB'),('RD','64GB')]
    # 确定数据库配置
    if mode == 'test':
        db_config = {
        'host': 'localhost',
        'user': 'remoteuser',
        'password': 'password',
        'database': 'cmsalpha',
        'charset': 'utf8mb4',
        'port': 3306,
        }
        reportDir = r'D:\Sync\业务报告\1on1\M-Test品质现况 202508.pptx'
        target_base_dir = r'D:\Sync\业务报告\1on1'
    else:
        db_config = {
            'host': '172.27.154.57',
            'user': 'remoteuser',
            'password': 'password',
            'database': 'cmsalpha',
            'charset': 'utf8mb4',
            'port': 3306,
        }
        reportDir = r'\\172.27.7.188\Mod_TestE\21. Fail Status Density\不良Status7.pptx'
        target_base_dir = r'\\172.27.7.188\Mod_TestE\21. Fail Status Density'
    # 确定db_fail_status最大月份
    max_month = get_max_month(db_config)

    # 获取当前系统月份的前一个月
    current_date = datetime.now()
    end_month = (current_date.replace(day=1) - timedelta(days=1)).strftime("%Y%m")

    # 生成当前年开始往前两年的年份列表，升序排列
    years = [current_date.year - i for i in range(3)]
    years.sort()
    # 生成上个月开始往前12个月的月份列表yyyymm格式，升序排列
    months = [
        (current_date.year, current_date.month - i) if current_date.month - i > 0 else (current_date.year - 1, 12 + (current_date.month - i))
        for i in range(1, 13)
    ]

    # max_month = '202412'

    # 比较 end_month 和 max_month，并逐月执行
    if end_month > max_month:
        # 从 max_month 的下一个月份开始
        current_month = increment_month(int(max_month))
        # 计算db_fail_status_density中需要的数值所有产品各个月份的结果
        while current_month <= int(end_month):
            month_str = str(current_month)
            result_df = process_data_from_db(db_config, month_str)
            if result_df is not None:
                print(f"Results for {month_str}:")
                print(result_df)
                # 将结果写入数据库
                Insert_data(result_df, db_config)
            current_month = increment_month(int(current_month))


    else:
        print("No new months to process.")
    for product in products:
        print(product)
        # 生成报表数据
        result = generate_report_failStatus(db_config, current_date, product)
        result = generate_target_year(db_config, result, product)
        print("Report data:", result)
        write_to_ppt(result, reportDir, product, 1)
        print("Write Complete")
        # # 写入PPT
        # list_shapes_in_slide(reportDir, 0)
        # # 复制文件到指定位置
        report_month = datetime.now().strftime('%Y%m')
        copy_and_save_report(reportDir, target_base_dir, report_month)
        # 增加月份


# Main function to execute processing
if __name__ == "__main__":
    main('test')